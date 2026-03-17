import asyncio
from pathlib import Path
from typing import IO, Any, List, Optional, Union
from uuid import uuid4

from agno.knowledge.chunking.document import DocumentChunking
from agno.knowledge.chunking.strategy import ChunkingStrategy, ChunkingStrategyType
from agno.knowledge.document.base import Document
from agno.knowledge.reader.base import Reader
from agno.knowledge.types import ContentType
from agno.utils.log import log_debug, log_error

try:
    from pptx import Presentation  # type: ignore
except ImportError:
    raise ImportError("The `python-pptx` package is not installed. Please install it via `pip install python-pptx`.")


class PPTXReader(Reader):
    """Reader for PPTX files"""

    def __init__(self, chunking_strategy: Optional[ChunkingStrategy] = DocumentChunking(), **kwargs):
        super().__init__(chunking_strategy=chunking_strategy, **kwargs)

    @classmethod
    def get_supported_chunking_strategies(self) -> List[ChunkingStrategyType]:
        """Get the list of supported chunking strategies for PPTX readers."""
        return [
            ChunkingStrategyType.DOCUMENT_CHUNKER,
            ChunkingStrategyType.CODE_CHUNKER,
            ChunkingStrategyType.FIXED_SIZE_CHUNKER,
            ChunkingStrategyType.SEMANTIC_CHUNKER,
            ChunkingStrategyType.AGENTIC_CHUNKER,
            ChunkingStrategyType.RECURSIVE_CHUNKER,
        ]

    @classmethod
    def get_supported_content_types(self) -> List[ContentType]:
        return [ContentType.PPTX]

    def read(self, file: Union[Path, IO[Any]], name: Optional[str] = None) -> List[Document]:
        """Read a pptx file and return a list of documents"""
        try:
            if isinstance(file, Path):
                if not file.exists():
                    raise FileNotFoundError(f"Could not find file: {file}")
                log_debug(f"Reading: {file}")
                presentation = Presentation(str(file))
                doc_name = name or file.stem
            else:
                log_debug(f"Reading uploaded file: {getattr(file, 'name', 'BytesIO')}")
                presentation = Presentation(file)
                doc_name = name or getattr(file, "name", "pptx_file").split(".")[0]

            # Extract text from all slides
            slide_texts = []
            for slide_number, slide in enumerate(presentation.slides, 1):
                slide_text = f"Slide {slide_number}:\n"

                # Extract text from shapes that contain text
                text_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text.strip())

                if text_content:
                    slide_text += "\n".join(text_content)
                else:
                    slide_text += "(No text content)"

                slide_texts.append(slide_text)

            doc_content = "\n\n".join(slide_texts)

            documents = [
                Document(
                    name=doc_name,
                    id=str(uuid4()),
                    content=doc_content,
                )
            ]

            if self.chunk:
                chunked_documents = []
                for document in documents:
                    chunked_documents.extend(self.chunk_document(document))
                return chunked_documents
            return documents

        except Exception as e:
            log_error(f"Error reading file: {e}")
            return []

    async def async_read(self, file: Union[Path, IO[Any]], name: Optional[str] = None) -> List[Document]:
        """Asynchronously read a pptx file and return a list of documents"""
        try:
            return await asyncio.to_thread(self.read, file, name)
        except Exception as e:
            log_error(f"Error reading file asynchronously: {e}")
            return []
